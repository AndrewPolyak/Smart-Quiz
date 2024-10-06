from pptx import Presentation
import docx
from PyPDF2 import PdfReader
import os
from tempfile import NamedTemporaryFile

class InputParser:
    def __init__(self):
        None

    
    def parse_ppt(self, powerpoint_path):
        '''
        Defines a method called "parse_ppt".
        '''
        ppt = Presentation(powerpoint_path) 

        text = ""
        
        # write text from powerpoint file into .txt file 
        for slide in ppt.slides:  
            for shape in slide.shapes:  
                if not shape.has_text_frame:  
                    continue 
                for paragraph in shape.text_frame.paragraphs:  
                    for run in paragraph.runs:  
                        text += run.text
        
        return text


    def parse_word(self, word_path):
        doc = docx.Document(word_path)

        text = ""

        for para in doc.paragraphs:
            text += para.text

        return text
    

    def parse_pdf(self, pdf_path):
        '''
        Defines a method called "parse__pdf".
        '''
        reader = PdfReader(pdf_path)

        text = ""

        for page in reader.pages:
            text += page.extract_text()

        return text


def process_file(file):
    text = ""
    file_name, extension = os.path.splitext(file.filename)
    extension = extension[1:].lower()  # Remove leading dot and make lowercase

    parser = InputParser()

    with NamedTemporaryFile(delete=False, suffix=extension) as temp_file:
        file.save(temp_file.name)
        temp_file_path = temp_file.name

    if extension == "pdf":
        text = parser.parse_pdf(temp_file_path)
    elif extension == "docx":
        text = parser.parse_word(temp_file_path)
    elif extension == "pptx":
        text = parser.parse_ppt(temp_file_path)

    os.remove(temp_file_path)

    return text
        
